#!/usr/bin/env python3
"""Fallback for Demonstration 1 if the live Codex build stalls. Synthetic only."""
from datetime import datetime, timedelta
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from openpyxl.styles import Font
from pptx import Presentation

OUT = Path(__file__).parent / "demo1_out"
OUT.mkdir(exist_ok=True)


def letter():
    d = Document()
    d.add_heading("Family update (teaching example)", 0)
    d.add_paragraph(
        "This note is a synthetic teaching example. It is not about a real child "
        "and it is not medical advice."
    )
    d.add_paragraph(
        "Your baby is in the hospital for RSV bronchiolitis. Overnight the oxygen "
        "level dipped and then improved on a small amount of oxygen by nasal cannula. "
        "Feeds are about half of usual. The team is watching breathing, oxygen, and feeding."
    )
    d.add_paragraph("This is a teaching example.")
    d.save(OUT / "family_update.docx")


def tracker():
    wb = Workbook()
    ws = wb.active
    ws.title = "wean"
    headers = ["datetime", "SpO2", "NC_L_min", "work_of_breathing_0_3", "feeds_pct", "hours_since_last_wean_attempt", "comments"]
    ws.append(headers)
    for cell in ws[1]:
        cell.font = Font(bold=True)
    start = datetime(2026, 1, 15, 18, 0)
    rows = [
        (0, 88, 0.0, 2, 40, "synthetic nadir on RA"),
        (4, 90, 0.5, 2, 40, "started 0.5 L"),
        (8, 92, 0.5, 1, 45, ""),
        (12, 94, 0.5, 1, 50, ""),
        (16, 94, 0.5, 1, 50, ""),
        (20, 95, 0.5, 1, 50, "still synthetic"),
    ]
    for i, (h, spo2, nc, wob, feeds, comment) in enumerate(rows, start=2):
        ws.cell(i, 1, start + timedelta(hours=h))
        ws.cell(i, 2, spo2)
        ws.cell(i, 3, nc)
        ws.cell(i, 4, wob)
        ws.cell(i, 5, feeds)
        if i == 2:
            ws.cell(i, 6, 0)
        else:
            ws.cell(i, 6, f"=ROUND((A{i}-A{i-1})*24,1)")
        ws.cell(i, 7, comment)
    ws["A10"] = "NOT FOR CLINICAL USE"
    wb.save(OUT / "oxygen_wean_tracker.xlsx")


def slides():
    prs = Presentation()
    titles = [
        "Synthetic bronchiolitis noon conference",
        "Objectives (teaching)",
        "Pathophysiology, cartoon-level",
        "Evidence bullets — citations not yet verified",
        "Wean principles (no doses)",
        "Discharge criteria (placeholder)",
        "What to remember",
        "References: citations not yet verified",
    ]
    for title in titles:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.text = "SYNTHETIC TEACHING CASE — NOT A REAL PATIENT"
        if "Evidence" in title or "References" in title:
            body.add_paragraph().text = "PLACEHOLDER—VERIFY BEFORE TEACHING"
    prs.save(OUT / "noon_conference.pptx")


if __name__ == "__main__":
    letter()
    tracker()
    slides()
    print("wrote", list(OUT.iterdir()))
